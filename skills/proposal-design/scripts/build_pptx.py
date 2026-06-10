#!/usr/bin/env python3
"""
可选：生成 pptx 的“骨架”文件（占位页），用于快速起稿。

依赖：
  pip install python-pptx

用法：
  python build_pptx.py <output.pptx>
"""

import sys


def main() -> int:
    if len(sys.argv) < 2:
        print("用法: python build_pptx.py <output.pptx>")
        return 2

    out_pptx = sys.argv[1]

    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:
        print(f"错误: 缺少依赖 python-pptx: {e}", file=sys.stderr)
        return 1

    prs = Presentation()

    titles = [
        "标题页",
        "我们怎么赢（1页）",
        "项目理解与范围边界",
        "叙事线与参观动线",
        "空间与视觉（多风格切换）",
        "系统架构（内容平台 + 中控联动）",
        "AI交互体验（至少4个）",
        "内容生产与更新机制",
        "实施计划与里程碑",
        "团队与保障",
        "评分点响应与证据",
        "结束页 / Q&A",
    ]

    layout = prs.slide_layouts[1]  # Title and Content
    for t in titles:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = t
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "（在此填充关键图/表/要点）"

    prs.save(out_pptx)
    print(f"已生成: {out_pptx}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

